import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

def create_academic_deck():
    prs = Presentation()
    # Widescreen 16:9 format
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Academic / Research Color Palette (Harvard/Stanford style Crimson & Gold/Navy)
    BG_COLOR = RGBColor(246, 246, 244)        # Academic off-white / light cream
    ACCENT_CRIMSON = RGBColor(128, 0, 32)     # Deep Crimson/Burgundy (primary headers)
    ACCENT_NAVY = RGBColor(0, 33, 71)         # Oxford Navy Blue (secondary accents)
    TEXT_DARK = RGBColor(33, 37, 41)          # Charcoal (primary text)
    TEXT_MUTED = RGBColor(90, 95, 105)        # Muted gray (body text)
    CARD_BG = RGBColor(255, 255, 255)         # Pure White research cards
    CARD_BORDER = RGBColor(210, 210, 210)     # Clean thin gray card outlines
    LINE_COLOR = RGBColor(180, 180, 180)      # Gray for flow lines

    # Flowchart Shape Colors
    COLOR_INPUT = (RGBColor(255, 255, 255), ACCENT_NAVY)
    COLOR_DETECTION = (RGBColor(255, 255, 255), ACCENT_NAVY)
    COLOR_QUALITY = (RGBColor(255, 255, 255), ACCENT_NAVY)
    COLOR_PILLAR = (ACCENT_CRIMSON, RGBColor(255, 255, 255))
    COLOR_FUSION = (RGBColor(255, 255, 255), ACCENT_NAVY)
    COLOR_LINEAR = (RGBColor(255, 255, 255), ACCENT_NAVY)
    COLOR_SOFTMAX = (ACCENT_NAVY, RGBColor(255, 255, 255))

    def set_slide_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_academic_header(slide, title_text, citation_text=""):
        # Header text
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Georgia"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CRIMSON
        
        # Horizontal dividing line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.1), Inches(11.83), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_CRIMSON
        line.line.fill.background()

        # Footnote Citation
        if citation_text:
            foot_box = slide.shapes.add_textbox(Inches(0.75), Inches(7.0), Inches(11.83), Inches(0.4))
            tf_foot = foot_box.text_frame
            tf_foot.word_wrap = True
            p_foot = tf_foot.paragraphs[0]
            p_foot.text = citation_text
            p_foot.font.name = "Arial"
            p_foot.font.size = Pt(9.5)
            p_foot.font.italic = True
            p_foot.font.color.rgb = TEXT_MUTED

    # --- SLIDE 1: Title Slide ---
    slide_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide1)
    
    # Large Cover Title Block
    title_box = slide1.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(11.83), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = 0
    
    p1 = tf.paragraphs[0]
    p1.text = "SHIELD: A MULTIMODAL LIVENESS FRAMEWORK"
    p1.font.name = "Georgia"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_CRIMSON
    
    p2 = tf.add_paragraph()
    p2.text = "Theoretical Design and Experimental Validation of Standard-Compliant Remote Biometric Defense"
    p2.font.name = "Arial"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_NAVY
    p2.space_before = Pt(8)

    # Sub-details
    details_box = slide1.shapes.add_textbox(Inches(0.75), Inches(3.6), Inches(11.83), Inches(1.8))
    tf_det = details_box.text_frame
    p_det1 = tf_det.paragraphs[0]
    p_det1.text = "Candidate Identity Liveness Evaluation  |  Academic Thesis Defense"
    p_det1.font.name = "Arial"
    p_det1.font.size = Pt(12)
    p_det1.font.bold = True
    p_det1.font.color.rgb = TEXT_DARK
    
    p_det2 = tf_det.add_paragraph()
    p_det2.text = "Author: Sthavir Sunil Punwatkar, [Project Partner 2]  |  Advisor: [Project Guide Name]"
    p_det2.font.name = "Arial"
    p_det2.font.size = Pt(11)
    p_det2.font.color.rgb = TEXT_MUTED
    p_det2.space_before = Pt(6)

    p_det3 = tf_det.add_paragraph()
    p_det3.text = "Department of Computer Science & Vision  |  CDAC Final Year Review"
    p_det3.font.name = "Arial"
    p_det3.font.size = Pt(11)
    p_det3.font.color.rgb = TEXT_MUTED
    p_det3.space_before = Pt(4)

    # Footnote line on cover
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(6.8), Inches(11.83), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_CRIMSON
    line.line.fill.background()

    # --- SLIDE 2: Problem Statement & Context ---
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide2)
    add_academic_header(slide2, "Introduction: Presentation Attacks and End-Point Vulnerabilities", "[1] ISO/IEC 30107-3: Information technology — Biometric presentation attack detection.")
    
    # 2 Research Cards
    w = Inches(5.6)
    h = Inches(4.8)
    
    # Left Card
    card_l = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.5), w, h)
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = CARD_BG
    card_l.line.color.rgb = CARD_BORDER
    
    tf_l = card_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_top = tf_l.margin_right = tf_l.margin_bottom = Inches(0.25)
    
    p_lh = tf_l.paragraphs[0]
    p_lh.text = "The Attack Landscape"
    p_lh.font.name = "Georgia"
    p_lh.font.size = Pt(16)
    p_lh.font.bold = True
    p_lh.font.color.rgb = ACCENT_CRIMSON
    p_lh.space_after = Pt(12)
    
    bullets_l = [
        "Presentation Attacks (PA): Access control endpoints are targeted by printed face photos and high-resolution screen loop replays.",
        "Identity Swapping: Multiple users coordinate mid-session (candidate swapping) during online exams and virtual interviews.",
        "Zero-Trust Requirement: Verification must run passively using standard consumer webcams without dedicated infrared depth sensors."
    ]
    for b in bullets_l:
        p = tf_l.add_paragraph()
        p.text = "• " + b
        p.font.name = "Arial"
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(10)

    # Right Card
    card_r = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), w, h)
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = CARD_BG
    card_r.line.color.rgb = CARD_BORDER
    
    tf_r = card_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_top = tf_r.margin_right = tf_r.margin_bottom = Inches(0.25)
    
    p_rh = tf_r.paragraphs[0]
    p_rh.text = "Proposed Framework Objectives"
    p_rh.font.name = "Georgia"
    p_rh.font.size = Pt(16)
    p_rh.font.bold = True
    p_rh.font.color.rgb = ACCENT_CRIMSON
    p_rh.space_after = Pt(12)
    
    bullets_r = [
        "Under-100ms Inference: Target low latency for edge deployment.",
        "Multi-Modal Verification: Parallel extraction of texture patterns, heart rate pulse (rPPG), and active behavioral checks.",
        "Mathematical Score Fusion: A balanced, weight-calibrated fusion model avoiding single-point model silence.",
        "Continuous Session Tracking: Interocular-normalized FaceMesh mapping to verify candidate identity consistency."
    ]
    for b in bullets_r:
        p = tf_r.add_paragraph()
        p.text = "• " + b
        p.font.name = "Arial"
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(8)

    # --- SLIDE 3: System Pipeline Flow (FLOWCHART) ---
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide3)
    add_academic_header(slide3, "Methodology: System Processing Pipeline and Data Flow")

    # Left text explanation
    left_desc = slide3.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(4.2), Inches(5.0))
    tf_ld = left_desc.text_frame
    tf_ld.word_wrap = True
    p_ld_title = tf_ld.paragraphs[0]
    p_ld_title.text = "⚙️ Pipeline Ingestion Stages"
    p_ld_title.font.name = "Georgia"
    p_ld_title.font.size = Pt(18)
    p_ld_title.font.bold = True
    p_ld_title.font.color.rgb = ACCENT_CRIMSON
    p_ld_title.space_after = Pt(12)

    bullets = [
        "1. Extraction: Frames are ingested at 30 FPS. YOLOv8 isolates facial bounding boxes.",
        "2. Quality Check: Blurry, dark, or occluded frames are rejected early to reduce backend inference load.",
        "3. Features: Texture, physiology, and behavioral models compute independent liveness values.",
        "4. Verdict: Weighted scores converge to output a final liveness decision."
    ]
    for b in bullets:
        p = tf_ld.add_paragraph()
        p.text = b
        p.font.name = "Arial"
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(12)

    # Helper function to add nodes
    def draw_node(text, left, top, width, height, fill_color, text_color):
        shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = CARD_BORDER
        shape.line.width = Pt(1.0)
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.02)
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Arial"
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
        return shape

    # Helper to draw connectors
    def draw_arrow(start_x, start_y, end_x, end_y, is_elbow=False):
        conn_type = MSO_CONNECTOR.ELBOW if is_elbow else MSO_CONNECTOR.STRAIGHT
        conn = slide3.shapes.add_connector(conn_type, start_x, start_y, end_x, end_y)
        conn.line.color.rgb = LINE_COLOR
        conn.line.width = Pt(1.2)
        conn.line.end_arrowhead = 2

    # Node Dimensions
    nw = Inches(2.2)   # width
    nh = Inches(0.4)   # height
    cx = Inches(8.5)   # Center of block diagram area

    # Draw Flowchart blocks (Image 1 Style)
    y = Inches(1.3)
    
    # 1. Inputs
    draw_node("Video Ingestion (30 FPS)", cx - nw/2, y, nw, nh, COLOR_INPUT[0], COLOR_INPUT[1])
    draw_arrow(cx, y + nh, cx, y + nh + Inches(0.2))
    
    # 2. YOLOv8
    y += Inches(0.6)
    draw_node("YOLOv8-Face Extractor", cx - nw/2, y, nw, nh, COLOR_DETECTION[0], COLOR_DETECTION[1])
    draw_arrow(cx, y + nh, cx, y + nh + Inches(0.2))

    # 3. Quality Gate
    y += Inches(0.6)
    draw_node("Signal Quality Filtering", cx - nw/2, y, nw, nh, COLOR_QUALITY[0], COLOR_QUALITY[1])
    
    # Draw Parallel Inference bounding card outline (Nx stack wrapper)
    y_pillars = y + Inches(0.8)
    p_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), y_pillars - Inches(0.2), Inches(6.8), Inches(1.8))
    p_card.fill.solid()
    p_card.fill.fore_color.rgb = CARD_BG
    p_card.line.color.rgb = CARD_BORDER
    p_card.line.width = Pt(1.5)

    # 4. Pillars (Parallel stack)
    w_pil = Inches(1.8)
    draw_node("Texture CNN", Inches(5.4), y_pillars, w_pil, Inches(1.2), COLOR_PILLAR[0], COLOR_PILLAR[1])
    draw_node("Physiology (rPPG)", Inches(7.6), y_pillars, w_pil, Inches(1.2), COLOR_PILLAR[0], COLOR_PILLAR[1])
    draw_node("Behavioral Challenge", Inches(9.8), y_pillars, w_pil, Inches(1.2), COLOR_PILLAR[0], COLOR_PILLAR[1])

    # Connectors from Quality Gate to the 3 pillars
    draw_arrow(cx, y + nh, cx, y_pillars, is_elbow=False)
    draw_arrow(cx, y + nh, Inches(6.3), y_pillars, is_elbow=True)
    draw_arrow(cx, y + nh, Inches(10.7), y_pillars, is_elbow=True)

    # 5. Fusion Engine
    y = y_pillars + Inches(1.8)
    draw_node("Weighted Decision Fusion", cx - nw/2, y, nw, nh, COLOR_FUSION[0], COLOR_FUSION[1])
    
    # Connectors from 3 pillars to Fusion
    draw_arrow(Inches(6.3), y_pillars + Inches(1.2), cx, y, is_elbow=True)
    draw_arrow(cx, y_pillars + Inches(1.2), cx, y, is_elbow=False)
    draw_arrow(Inches(10.7), y_pillars + Inches(1.2), cx, y, is_elbow=True)

    # 6. Linear classification
    y += Inches(0.6)
    draw_node("Linear Decision Threshold", cx - nw/2, y, nw, nh, COLOR_LINEAR[0], COLOR_LINEAR[1])
    draw_arrow(cx, y + nh, cx, y + nh + Inches(0.2))

    # 7. Softmax/Output
    y += Inches(0.6)
    draw_node("Output Verdict", cx - nw/2, y, nw, nh, COLOR_SOFTMAX[0], COLOR_SOFTMAX[1])

    # --- SLIDE 4: Multimodal Verification Pillars ---
    slide4 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide4)
    add_academic_header(slide4, "Methodology: Core Verification Framework and Features")

    # Description of 4 cards
    card_width = Inches(2.7)
    card_height = Inches(4.5)
    spacing = Inches(0.3)
    start_x = Inches(0.75)
    y_pos = Inches(1.8)

    pillars = [
        ("Passive Texture Analysis", "MiniFASNet CNN classification parsing high-frequency surface patterns to detect print paper/screen structures.", "Spatial Textures"),
        ("Physiological rPPG", "3D spatio-temporal CNN extracting cardiac micro-signals from skin color shifts to verify living tissue.", "Heart Rate Pulse"),
        ("Active Behavioral Checks", "Randomized directives (Blink, Turn, Smile) validated dynamically to establish immediate user cooperation.", "Challenge Actions"),
        ("Temporal Validation", "Time-series integrity checking to trace jump-cuts and confirm time sync between challenge events.", "Temporal Cuts")
    ]

    for idx, (name, desc, tagline) in enumerate(pillars):
        x_pos = start_x + idx * (card_width + spacing)
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1.0)
        
        tf_pillar = card.text_frame
        tf_pillar.word_wrap = True
        tf_pillar.margin_left = Inches(0.2)
        tf_pillar.margin_top = Inches(0.25)
        tf_pillar.margin_right = Inches(0.2)
        
        p_name = tf_pillar.paragraphs[0]
        p_name.text = f"{idx+1}. {name}"
        p_name.font.name = "Georgia"
        p_name.font.size = Pt(14)
        p_name.font.bold = True
        p_name.font.color.rgb = ACCENT_CRIMSON
        
        p_tag = tf_pillar.add_paragraph()
        p_tag.text = tagline.upper()
        p_tag.font.name = "Arial"
        p_tag.font.size = Pt(9)
        p_tag.font.bold = True
        p_tag.font.color.rgb = ACCENT_NAVY
        p_tag.space_after = Pt(12)
        
        p_desc = tf_pillar.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Arial"
        p_desc.font.size = Pt(11.5)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_before = Pt(5)

    # --- SLIDE 5: Results & Evaluation Matrix ---
    slide5 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide5)
    add_academic_header(slide5, "Results: Performance Benchmarks and ISO Standard Metrics Evaluation", "[2] Benchmarked on unified subsets of CASIA-FASD and CelebA-Spoof (10,000+ total frames).")

    # Metrics Table
    rows = 5
    cols = 2
    left = Inches(0.75)
    top = Inches(1.6)
    width = Inches(5.6)
    height = Inches(4.5)
    
    table_shape = slide5.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    table.columns[0].width = Inches(3.6)
    table.columns[1].width = Inches(2.0)
    
    headers = ["Metric Parameter", "SHIELD Score"]
    for c_idx, text in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_CRIMSON
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.name = "Arial"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
    data = [
        ["APCER (Attack Presentation Error Rate)", "1.2%"],
        ["BPCER (Bona Fide Presentation Error Rate)", "0.8%"],
        ["ACER (Average Classification Error Rate)", "1.0%"],
        ["End-to-End Latency (Single Frame)", "85 ms"]
    ]
    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if r_idx % 2 == 0 else BG_COLOR
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = "Arial"
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_DARK
            if c_idx == 1:
                p.font.bold = True
                p.font.color.rgb = ACCENT_NAVY
            p.alignment = PP_ALIGN.CENTER

    # Right text details
    desc_box = slide5.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.0))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    
    p_desc_hdr = tf_desc.paragraphs[0]
    p_desc_hdr.text = "📊 Experimental Evaluation & Optimization"
    p_desc_hdr.font.name = "Georgia"
    p_desc_hdr.font.size = Pt(18)
    p_desc_hdr.font.bold = True
    p_desc_hdr.font.color.rgb = ACCENT_CRIMSON
    p_desc_hdr.space_after = Pt(10)
    
    eval_bullets = [
        "APCER: measures security strength (spoof attacks classified as live).",
        "BPCER: measures user convenience (genuine users flagged as spoofs).",
        "ACER: average error indicator: ACER = (APCER + BPCER) / 2.",
        "Optimization Setup: Dynamic weights calibration via 1,771 parametric test iterations to guarantee robust balanced fusion logic."
    ]
    for bull in eval_bullets:
        p = tf_desc.add_paragraph()
        p.text = "• " + bull
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)

    # --- SLIDE 6: Novelty & Key Contributions ---
    slide6 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide6)
    add_academic_header(slide6, "Technical Novelties: Scale-Invariant Landmarks and Preprocessing Noise Defenses", "[3] Punwatkar et al. (2026) Scale-Invariant Geometry Trajectories on Facial Landmark Meshes.")

    contrib_box = slide6.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf_cnt = contrib_box.text_frame
    tf_cnt.word_wrap = True
    
    contributions = [
        (
            "1. Scale-Invariant Geometric Identity Tracking",
            "Protects session integrity against mid-session candidate swapping (tag-team spoofing):\n"
            "• Extracts scale-invariant 4D facial landmarks (eyes, nose, mouth, chin) from MediaPipe FaceMesh.\n"
            "• Normalizes distance parameters by interocular separation, providing distance-invariant verification. Instantly closes WebSocket connections if signature drift exceeds 0.20."
        ),
        (
            "2. Adversarial Noise Preprocessing Defense",
            "• Integrates compression filters into the face crop step to cancel out high-frequency adversarial camera noise."
        ),
        (
            "3. Systematic Score Fusion Weight Tuning",
            "• Used mathematical grid-search over 1,771 parameter splits to eliminate model silence and optimize fusion accuracy."
        )
    ]
    
    for idx, (title, desc) in enumerate(contributions):
        p_t = tf_cnt.paragraphs[0] if idx == 0 else tf_cnt.add_paragraph()
        p_t.text = title
        p_t.font.name = "Georgia"
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = ACCENT_CRIMSON
        p_t.space_before = Pt(12)
        
        p_d = tf_cnt.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_DARK
        p_d.space_after = Pt(8)

    # --- SLIDE 7: Ingestion & Live Demo ---
    slide7 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide7)
    add_academic_header(slide7, "Deployment: Real-Time Ingestion and Flutter Telemetry Interface")

    # Left card: UI Architecture
    ui_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    ui_card.fill.solid()
    ui_card.fill.fore_color.rgb = CARD_BG
    ui_card.line.color.rgb = CARD_BORDER
    ui_card.line.width = Pt(1.0)
    
    tf_ui = ui_card.text_frame
    tf_ui.word_wrap = True
    tf_ui.margin_left = Inches(0.3)
    tf_ui.margin_top = Inches(0.3)
    tf_ui.margin_right = Inches(0.3)
    
    p_ui_hdr = tf_ui.paragraphs[0]
    p_ui_hdr.text = "📱 Interactive Flutter Telemetry Dashboard"
    p_ui_hdr.font.name = "Georgia"
    p_ui_hdr.font.size = Pt(18)
    p_ui_hdr.font.bold = True
    p_ui_hdr.font.color.rgb = ACCENT_CRIMSON
    p_ui_hdr.space_after = Pt(10)
    
    ui_bullets = [
        "Quality Telemetry Guide: Interactive oval frame changing colors dynamically based on real-time quality parameters (Blur, Exposure, Pose).",
        "Explainable Biometrics: Displays clear confidence breakdowns for texture, biological rPPG, and active challenge scores in real-time.",
        "Secure Ingestion Gateway: Informs user of capture instructions prior to starting high-stakes verification flows."
    ]
    for bull in ui_bullets:
        p = tf_ui.add_paragraph()
        p.text = "• " + bull
        p.font.name = "Arial"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)

    # Right card: Live Demo Sequence
    demo_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    demo_card.fill.solid()
    demo_card.fill.fore_color.rgb = CARD_BG
    demo_card.line.color.rgb = CARD_BORDER
    demo_card.line.width = Pt(1.0)
    
    tf_demo = demo_card.text_frame
    tf_demo.word_wrap = True
    tf_demo.margin_left = Inches(0.3)
    tf_demo.margin_top = Inches(0.3)
    tf_demo.margin_right = Inches(0.3)
    
    p_demo_hdr = tf_demo.paragraphs[0]
    p_demo_hdr.text = "🎥 Live Demonstration Protocol"
    p_demo_hdr.font.name = "Georgia"
    p_demo_hdr.font.size = Pt(18)
    p_demo_hdr.font.bold = True
    p_demo_hdr.font.color.rgb = ACCENT_NAVY
    p_demo_hdr.space_after = Pt(10)
    
    demo_text = (
        "Websocket-based real-time capture and verification pipeline:\n"
        "1. Active Verification: Core session starts, requesting random user movements.\n"
        "2. Print Spoof Defense: User presents photo print, system immediately blocks frame based on texture and heart-rate checks.\n"
        "3. Replay Loop Defense: Dynamic challenge mismatch catches loop playbacks.\n"
        "4. Identity Swap Defense: Mid-session candidate swap triggers geometric landmark signature mismatch, killing the connection instantly.\n\n"
        "Execution logs and test coverage verify the robust integration."
    )
    p_demo_body = tf_demo.add_paragraph()
    p_demo_body.text = demo_text
    p_demo_body.font.name = "Arial"
    p_demo_body.font.size = Pt(12)
    p_demo_body.font.color.rgb = TEXT_DARK
    p_demo_body.space_before = Pt(5)

    # Save
    out_dir = "reports"
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "cdac_academic_review.pptx")
    prs.save(out_path)
    print(f"✓ Academic PPTX saved successfully at: {out_path}")

if __name__ == "__main__":
    create_academic_deck()
