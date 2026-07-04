import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

def create_deck():
    prs = Presentation()
    # Widescreen 16:9 format
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Custom Pitch Deck Color Palette (Matching Uploaded Media 0)
    BG_COLOR = RGBColor(10, 10, 16)          # Ultra dark charcoal/black
    ACCENT_PURPLE = RGBColor(138, 43, 226)   # Electric violet (accent highlights)
    TEXT_WHITE = RGBColor(255, 255, 255)     # Clean white (primary text)
    TEXT_MUTED = RGBColor(200, 200, 220)     # Soft light gray/lavender (body text)
    TEXT_DARK_GRAY = RGBColor(130, 120, 150) # Muted slate (sub-info)
    CARD_BG = RGBColor(24, 22, 34)           # Muted card background
    CARD_BORDER = RGBColor(60, 50, 90)       # Dark violet border

    # Flowchart Shape Colors (Matching Uploaded Media 1)
    COLOR_INPUT = (RGBColor(248, 187, 208), RGBColor(40, 40, 40))   # Muted Pink
    COLOR_DETECTION = (RGBColor(255, 224, 178), RGBColor(40, 40, 40)) # Muted Orange
    COLOR_QUALITY = (RGBColor(255, 249, 196), RGBColor(40, 40, 40))   # Muted Yellow
    COLOR_PILLAR = (RGBColor(224, 242, 241), RGBColor(40, 40, 40))    # Muted Teal
    COLOR_FUSION = (RGBColor(255, 249, 196), RGBColor(40, 40, 40))    # Muted Yellow
    COLOR_LINEAR = (RGBColor(225, 190, 231), RGBColor(40, 40, 40))    # Muted Purple
    COLOR_SOFTMAX = (RGBColor(200, 230, 201), RGBColor(40, 40, 40))   # Muted Green

    # Image Paths
    SHIELD_COVER_IMG = "/home/sp/.gemini/antigravity-cli/brain/25018484-f1af-4d12-9a28-f0e9940aa64f/shield_cover_visual_1783185170999.jpg"
    TELEMETRY_IMG = "/home/sp/.gemini/antigravity-cli/brain/25018484-f1af-4d12-9a28-f0e9940aa64f/telemetry_visual_1783185198779.jpg"
    FACIAL_MESH_IMG = "/home/sp/.gemini/antigravity-cli/brain/25018484-f1af-4d12-9a28-f0e9940aa64f/facial_mesh_visual_1783185185503.jpg"

    def set_slide_bg(slide, image_path=None):
        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        else:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = BG_COLOR

    def add_title(slide, text, color=ACCENT_PURPLE, left=Inches(0.75), top=Inches(0.4), width=Inches(11.83)):
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = 0
        tf.margin_top = 0
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Arial"
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = color
        return title_box

    # --- SLIDE 1: Title and Problem Statement (Cover) ---
    slide_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide1, SHIELD_COVER_IMG)
    
    # Title details overlay
    title_box = slide1.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(5.5), Inches(1.6))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    tf.margin_right = 0
    tf.margin_bottom = 0
    
    p1 = tf.paragraphs[0]
    p1.text = "SHIELD"
    p1.font.name = "Arial"
    p1.font.size = Pt(42)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "Multimodal Real-Time Biometric Liveness & Identity Verification System"
    p2.font.name = "Arial"
    p2.font.size = Pt(15)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_PURPLE
    p2.space_before = Pt(6)

    # Project Metadata (Partners & Guide)
    meta_box = slide1.shapes.add_textbox(Inches(0.75), Inches(2.1), Inches(5.5), Inches(1.1))
    tf_meta = meta_box.text_frame
    tf_meta.word_wrap = True
    tf_meta.margin_left = 0
    tf_meta.margin_top = 0
    tf_meta.margin_right = 0
    tf_meta.margin_bottom = 0
    
    p_meta1 = tf_meta.paragraphs[0]
    p_meta1.text = "Project Partners: Sthavir Sunil Punwatkar, [Project Partner 2]"
    p_meta1.font.name = "Arial"
    p_meta1.font.size = Pt(11.5)
    p_meta1.font.bold = True
    p_meta1.font.color.rgb = TEXT_MUTED
    
    p_meta2 = tf_meta.add_paragraph()
    p_meta2.text = "Project Guide: [Project Guide Name]"
    p_meta2.font.name = "Arial"
    p_meta2.font.size = Pt(11.5)
    p_meta2.font.bold = True
    p_meta2.font.color.rgb = TEXT_MUTED
    p_meta2.space_before = Pt(3)

    p_meta3 = tf_meta.add_paragraph()
    p_meta3.text = "CDAC Project Review  |  Domain: Biometrics & Computer Vision"
    p_meta3.font.name = "Arial"
    p_meta3.font.size = Pt(11.5)
    p_meta3.font.color.rgb = TEXT_DARK_GRAY
    p_meta3.space_before = Pt(5)

    # Problem Statement Card (Left Overlay)
    card1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(3.3), Inches(5.5), Inches(3.6))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = CARD_BORDER
    card1.line.width = Pt(1.0)
    
    tf_card = card1.text_frame
    tf_card.word_wrap = True
    tf_card.margin_left = Inches(0.25)
    tf_card.margin_top = Inches(0.25)
    tf_card.margin_right = Inches(0.25)
    tf_card.margin_bottom = Inches(0.25)
    
    p_prob_hdr = tf_card.paragraphs[0]
    p_prob_hdr.text = "PROBLEM FORMULATION"
    p_prob_hdr.font.name = "Arial"
    p_prob_hdr.font.size = Pt(15)
    p_prob_hdr.font.bold = True
    p_prob_hdr.font.color.rgb = ACCENT_PURPLE
    
    p_prob_body = tf_card.add_paragraph()
    p_prob_body.text = (
        "Remote biometric systems face vulnerabilities to spoofing:\n"
        "• Presentation Attacks (PA): Static photo prints, video loop replays, and 3D masks deceive traditional facial recognition.\n"
        "• Identity Swapping: User swaps mid-session (tag-team cheating) during active testing or authentication.\n"
        "• Deployment Limits: Existing anti-spoofing methods require specialized depth hardware, rendering wide scale deployments impractical."
    )
    p_prob_body.font.name = "Arial"
    p_prob_body.font.size = Pt(11)
    p_prob_body.font.color.rgb = TEXT_MUTED
    p_prob_body.space_before = Pt(6)

    # --- SLIDE 2: Objectives (Standard Layout) ---
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide2)
    add_title(slide2, "Project Objectives")

    # Left Column: Primary Goals
    left_box = slide2.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.6), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p_l1 = tf_left.paragraphs[0]
    p_l1.text = "🎯 Core Objectives"
    p_l1.font.name = "Arial"
    p_l1.font.size = Pt(18)
    p_l1.font.bold = True
    p_l1.font.color.rgb = ACCENT_PURPLE
    p_l1.space_after = Pt(10)
    
    objectives_list = [
        "Develop an end-to-end, real-time liveness pipeline processing multi-layered verification streams in <100ms.",
        "Implement non-intrusive biological checks by extracting blood volume pulse (BVP) from standard RGB video feeds (rPPG).",
        "Introduce randomized active challenge-response tasks combined with frame validation to counter pre-recorded replays.",
        "Construct an explainable multi-modal decision fusion engine to facilitate detailed audit trails."
    ]
    for obj in objectives_list:
        p = tf_left.add_paragraph()
        p.text = "• " + obj
        p.font.name = "Arial"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(12)

    # Right Column: High-Level Targets
    right_box = slide2.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.0))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p_r1 = tf_right.paragraphs[0]
    p_r1.text = "⚙️ Engineering Targets & KPIs"
    p_r1.font.name = "Arial"
    p_r1.font.size = Pt(18)
    p_r1.font.bold = True
    p_r1.font.color.rgb = TEXT_WHITE
    p_r1.space_after = Pt(10)
    
    kpis = [
        ("ISO/IEC 30107-3 Standard Adherence", "Deliver an ACER (Average Classification Error Rate) under 1.5% to verify enterprise-level security."),
        ("Edge-Deployable Performance", "Optimize inference workloads to support execution in under 100ms on consumer-grade CPUs."),
        ("Zero-Dependency RGB Deployment", "Perform biological, texture, and active checks using simple 2D webcams without dedicated sensors.")
    ]
    for title, desc in kpis:
        p_t = tf_right.add_paragraph()
        p_t.text = "▪️ " + title
        p_t.font.name = "Arial"
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = ACCENT_PURPLE
        p_t.space_before = Pt(8)
        
        p_d = tf_right.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_after = Pt(6)

    # --- SLIDE 3: System Architecture & Flow (FLOWCHART BUILD) ---
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide3)
    add_title(slide3, "System Architecture Pipeline")

    # Left text explanation
    left_desc = slide3.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(4.2), Inches(5.0))
    tf_ld = left_desc.text_frame
    tf_ld.word_wrap = True
    p_ld_title = tf_ld.paragraphs[0]
    p_ld_title.text = "⚙️ Modular Data Pipeline"
    p_ld_title.font.name = "Arial"
    p_ld_title.font.size = Pt(18)
    p_ld_title.font.bold = True
    p_ld_title.font.color.rgb = ACCENT_PURPLE
    p_ld_title.space_after = Pt(12)

    bullets = [
        "A modular layer architecture styled after high-performance neural blocks (like Transformer stacks).",
        "Frames pass sequentially through detection, validation, and quality gates before running parallel inferences.",
        "Texture (Spatial), Physiology (Biological BVP), and Challenge (Temporal Actions) operate as distinct pipelines.",
        "Decision Fusion acts as the linear convergence layer to calculate output probabilities."
    ]
    for b in bullets:
        p = tf_ld.add_paragraph()
        p.text = "• " + b
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(10)

    # Helper function to add nodes
    def draw_node(text, left, top, width, height, fill_color, text_color):
        shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = CARD_BORDER
        shape.line.width = Pt(1.0)
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.02)
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Arial"
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
        return shape

    # Helper to draw connectors
    def draw_arrow(start_x, start_y, end_x, end_y, is_elbow=False):
        conn_type = MSO_CONNECTOR.ELBOW if is_elbow else MSO_CONNECTOR.STRAIGHT
        conn = slide3.shapes.add_connector(conn_type, start_x, start_y, end_x, end_y)
        conn.line.color.rgb = CARD_BORDER
        conn.line.width = Pt(1.2)
        conn.line.end_arrowhead = 2 # Triangle arrow

    # Node Dimensions
    nw = Inches(2.2)   # width
    nh = Inches(0.4)   # height
    cx = Inches(8.5)   # Center of block diagram area

    # Draw Flowchart blocks (Image 1 Style)
    y = Inches(1.1)
    
    # 1. Inputs
    draw_node("Video Frames (BGR stream)", cx - nw/2, y, nw, nh, COLOR_INPUT[0], COLOR_INPUT[1])
    draw_arrow(cx, y + nh, cx, y + nh + Inches(0.2))
    
    # 2. YOLOv8
    y += Inches(0.6)
    draw_node("YOLOv8-Face Detection", cx - nw/2, y, nw, nh, COLOR_DETECTION[0], COLOR_DETECTION[1])
    draw_arrow(cx, y + nh, cx, y + nh + Inches(0.2))

    # 3. Quality Gate
    y += Inches(0.6)
    draw_node("Signal Quality Filter Gate", cx - nw/2, y, nw, nh, COLOR_QUALITY[0], COLOR_QUALITY[1])
    
    # Draw Parallel Inference bounding card outline (Nx stack wrapper)
    y_pillars = y + Inches(0.8)
    p_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), y_pillars - Inches(0.2), Inches(6.8), Inches(1.8))
    p_card.fill.background()
    p_card.line.color.rgb = CARD_BORDER
    p_card.line.width = Pt(1.5)
    
    # Bounding card title (NX label)
    lbl = slide3.shapes.add_textbox(Inches(11.35), y_pillars + Inches(0.4), Inches(0.5), Inches(0.4))
    lbl.text_frame.paragraphs[0].text = "3x"
    lbl.text_frame.paragraphs[0].font.name = "Arial"
    lbl.text_frame.paragraphs[0].font.size = Pt(14)
    lbl.text_frame.paragraphs[0].font.bold = True
    lbl.text_frame.paragraphs[0].font.color.rgb = ACCENT_PURPLE

    # 4. Pillars (Parallel stack)
    w_pil = Inches(1.8)
    draw_node("Passive Texture CNN", Inches(5.4), y_pillars, w_pil, Inches(1.2), COLOR_PILLAR[0], COLOR_PILLAR[1])
    draw_node("Physiological rPPG", Inches(7.6), y_pillars, w_pil, Inches(1.2), COLOR_PILLAR[0], COLOR_PILLAR[1])
    draw_node("Active Challenges", Inches(9.8), y_pillars, w_pil, Inches(1.2), COLOR_PILLAR[0], COLOR_PILLAR[1])

    # Connectors from Quality Gate to the 3 pillars
    draw_arrow(cx, y + nh, cx, y_pillars, is_elbow=False)
    draw_arrow(cx, y + nh, Inches(6.3), y_pillars, is_elbow=True)
    draw_arrow(cx, y + nh, Inches(10.7), y_pillars, is_elbow=True)

    # 5. Fusion Engine
    y = y_pillars + Inches(1.8)
    draw_node("Weighted Decision Fusion", cx - nw/2, y, nw, nh, COLOR_FUSION[0], COLOR_FUSION[1])
    
    # Connectors from 3 pillars to Fusion
    draw_arrow(Inches(6.3), y_pillars + Inches(1.2), cx, y, is_elbow=True)
    draw_arrow(cx, y_pillars + Inches(1.2), cx, y, is_elbow=False)
    draw_arrow(Inches(10.7), y_pillars + Inches(1.2), cx, y, is_elbow=True)

    # 6. Linear classification
    y += Inches(0.6)
    draw_node("Linear Decision Threshold", cx - nw/2, y, nw, nh, COLOR_LINEAR[0], COLOR_LINEAR[1])
    draw_arrow(cx, y + nh, cx, y + nh + Inches(0.2))

    # 7. Softmax/Output
    y += Inches(0.6)
    draw_node("Output Verdict (Live/Spoof)", cx - nw/2, y, nw, nh, COLOR_SOFTMAX[0], COLOR_SOFTMAX[1])

    # --- SLIDE 4: Multimodal Verification Pillars (Standard Layout) ---
    slide4 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide4)
    add_title(slide4, "Multimodal Verification Pillars")

    # Description of 4 cards
    card_width = Inches(2.7)
    card_height = Inches(4.5)
    spacing = Inches(0.3)
    start_x = Inches(0.75)
    y_pos = Inches(1.8)

    pillars = [
        ("Passive Texture", "MiniFASNet CNN classification parsing high-frequency surface patterns to detect print paper/screen structures.", "Spatial Textures"),
        ("Physiological rPPG", "3D spatio-temporal CNN extracting cardiac micro-signals from skin color shifts to verify living tissue.", "Heart Rate Pulse"),
        ("Active Challenges", "Randomized directives (Blink, Turn, Smile) validated dynamically to establish immediate user cooperation.", "Challenge Actions"),
        ("Temporal Validator", "Time-series integrity checking to trace jump-cuts and confirm time sync between challenge events.", "Temporal Cuts")
    ]

    for idx, (name, desc, tagline) in enumerate(pillars):
        x_pos = start_x + idx * (card_width + spacing)
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1.0)
        
        tf_pillar = card.text_frame
        tf_pillar.word_wrap = True
        tf_pillar.margin_left = Inches(0.2)
        tf_pillar.margin_top = Inches(0.25)
        tf_pillar.margin_right = Inches(0.2)
        
        p_name = tf_pillar.paragraphs[0]
        p_name.text = f"{idx+1}. {name}"
        p_name.font.name = "Arial"
        p_name.font.size = Pt(15)
        p_name.font.bold = True
        p_name.font.color.rgb = TEXT_WHITE
        
        p_tag = tf_pillar.add_paragraph()
        p_tag.text = tagline.upper()
        p_tag.font.name = "Arial"
        p_tag.font.size = Pt(9)
        p_tag.font.bold = True
        p_tag.font.color.rgb = ACCENT_PURPLE
        p_tag.space_after = Pt(12)
        
        p_desc = tf_pillar.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Arial"
        p_desc.font.size = Pt(11.5)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_before = Pt(5)

    # --- SLIDE 5: Results & Evaluation Matrix (Image Backdrop Overlay) ---
    slide5 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide5, TELEMETRY_IMG)
    
    add_title(slide5, "Results & Evaluation Matrix", TEXT_WHITE, Inches(0.75), Inches(0.4), Inches(5.6))

    # Metrics Table
    rows = 5
    cols = 2
    left = Inches(0.75)
    top = Inches(1.3)
    width = Inches(5.6)
    height = Inches(2.2)
    
    table_shape = slide5.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    table.columns[0].width = Inches(3.6)
    table.columns[1].width = Inches(2.0)
    
    headers = ["Metric Parameter", "SHIELD Score"]
    for c_idx, text in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.name = "Arial"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = ACCENT_PURPLE
        p.alignment = PP_ALIGN.CENTER
        
    data = [
        ["APCER (Attack Error Rate)", "1.2%"],
        ["BPCER (Bona Fide Error Rate)", "0.8%"],
        ["ACER (Average Error Rate)", "1.0%"],
        ["End-to-End Inference Latency", "85 ms"]
    ]
    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_COLOR
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = "Arial"
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_MUTED
            if c_idx == 1:
                p.font.bold = True
                p.font.color.rgb = TEXT_WHITE
            p.alignment = PP_ALIGN.CENTER

    # Bounding Box for Details (Evaluation matrix definitions)
    desc_box = slide5.shapes.add_textbox(Inches(0.75), Inches(3.7), Inches(5.6), Inches(3.3))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    
    p_desc_hdr = tf_desc.paragraphs[0]
    p_desc_hdr.text = "📊 ISO/IEC 30107-3 Standard Metrics"
    p_desc_hdr.font.name = "Arial"
    p_desc_hdr.font.size = Pt(15)
    p_desc_hdr.font.bold = True
    p_desc_hdr.font.color.rgb = ACCENT_PURPLE
    p_desc_hdr.space_after = Pt(6)
    
    eval_bullets = [
        "APCER: Attack Presentation Classification Error Rate. Measures the % of spoof attacks incorrectly classified as live.",
        "BPCER: Bona Fide Presentation Classification Error Rate. Measures the % of genuine live users incorrectly flagged as spoofs.",
        "ACER: Average Classification Error Rate. Calculated as the average of APCER and BPCER: ACER = (APCER + BPCER) / 2.",
        "Empirical Calibration: System weights tuned via 1,771 test combinations under min_weight=0.10 constraint to achieve 1.0% ACER."
    ]
    for bull in eval_bullets:
        p = tf_desc.add_paragraph()
        p.text = "• " + bull
        p.font.name = "Arial"
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(4)

    # --- SLIDE 6: Novelty & Key Contributions (Image Backdrop Overlay) ---
    slide6 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide6, FACIAL_MESH_IMG)
    
    add_title(slide6, "Novelty & Custom Contributions", TEXT_WHITE, Inches(0.75), Inches(0.4), Inches(5.6))

    contrib_box = slide6.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(5.6), Inches(5.8))
    tf_cnt = contrib_box.text_frame
    tf_cnt.word_wrap = True
    
    contributions = [
        (
            "1. Real-Time Scale-Invariant Identity Consistency Check",
            "Defeats mid-session candidate swapping (tag-team fraud):\n"
            "• Computes scale-invariant 4D geometric signature ratios (nose, eyes, chin, lip corners) using MediaPipe FaceMesh landmarks.\n"
            "• Signature distance is verified frame-by-frame. Websocket is killed immediately if signature drift exceeds threshold (>0.20)."
        ),
        (
            "2. JPEG Compression Defenses",
            "Enhances resilience against adversarial spatial noise attacks:\n"
            "• Integrates compression filters into the face crop step to cancel out artificial noise patterns generated by digital camera overlays."
        ),
        (
            "3. Systematic Grid Search Optimization",
            "• Evaluated 1,771 separate sensor weight parameter combinations mathematically to eliminate sub-model silence."
        )
    ]
    
    for idx, (title, desc) in enumerate(contributions):
        p_t = tf_cnt.paragraphs[0] if idx == 0 else tf_cnt.add_paragraph()
        p_t.text = title
        p_t.font.name = "Arial"
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = ACCENT_PURPLE
        p_t.space_before = Pt(8)
        
        p_d = tf_cnt.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_after = Pt(6)

    # --- SLIDE 7: Demo Video & Live Telemetry UI (Standard Layout) ---
    slide7 = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide7)
    add_title(slide7, "Interface Design & Live Demonstration")

    # Left card: UI Architecture
    ui_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.5))
    ui_card.fill.solid()
    ui_card.fill.fore_color.rgb = CARD_BG
    ui_card.line.color.rgb = CARD_BORDER
    ui_card.line.width = Pt(1.0)
    
    tf_ui = ui_card.text_frame
    tf_ui.word_wrap = True
    tf_ui.margin_left = Inches(0.3)
    tf_ui.margin_top = Inches(0.3)
    tf_ui.margin_right = Inches(0.3)
    
    p_ui_hdr = tf_ui.paragraphs[0]
    p_ui_hdr.text = "📱 Interactive Flutter Telemetry Dashboard"
    p_ui_hdr.font.name = "Arial"
    p_ui_hdr.font.size = Pt(18)
    p_ui_hdr.font.bold = True
    p_ui_hdr.font.color.rgb = TEXT_WHITE
    p_ui_hdr.space_after = Pt(10)
    
    ui_bullets = [
        "Quality Telemetry Guide: Interactive oval frame changing colors dynamically based on real-time quality parameters (Blur, Exposure, Pose).",
        "Explainable Biometrics: Displays clear confidence breakdowns for texture, biological rPPG, and active challenge scores in real-time.",
        "Secure Ingestion Gateway: Informs user of capture instructions prior to starting high-stakes verification flows."
    ]
    for bull in ui_bullets:
        p = tf_ui.add_paragraph()
        p.text = "• " + bull
        p.font.name = "Arial"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(8)

    # Right card: Live Demo Sequence
    demo_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))
    demo_card.fill.solid()
    demo_card.fill.fore_color.rgb = CARD_BG
    demo_card.line.color.rgb = CARD_BORDER
    demo_card.line.width = Pt(1.0)
    
    tf_demo = demo_card.text_frame
    tf_demo.word_wrap = True
    tf_demo.margin_left = Inches(0.3)
    tf_demo.margin_top = Inches(0.3)
    tf_demo.margin_right = Inches(0.3)
    
    p_demo_hdr = tf_demo.paragraphs[0]
    p_demo_hdr.text = "🎥 Live Demonstration Protocol"
    p_demo_hdr.font.name = "Arial"
    p_demo_hdr.font.size = Pt(18)
    p_demo_hdr.font.bold = True
    p_demo_hdr.font.color.rgb = ACCENT_PURPLE
    p_demo_hdr.space_after = Pt(10)
    
    demo_text = (
        "Websocket-based real-time capture and verification pipeline:\n"
        "1. Active Verification: Core session starts, requesting random user movements.\n"
        "2. Print Spoof Defense: User presents photo print, system immediately blocks frame based on texture and heart-rate checks.\n"
        "3. Replay Loop Defense: Dynamic challenge mismatch catches loop playbacks.\n"
        "4. Identity Swap Defense: Mid-session candidate swap triggers geometric landmark signature mismatch, killing the connection instantly.\n\n"
        "Execution logs and test coverage verify the robust integration."
    )
    p_demo_body = tf_demo.add_paragraph()
    p_demo_body.text = demo_text
    p_demo_body.font.name = "Arial"
    p_demo_body.font.size = Pt(12)
    p_demo_body.font.color.rgb = TEXT_MUTED
    p_demo_body.space_before = Pt(5)

    # Save
    out_dir = "reports"
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "cdac_project_review.pptx")
    prs.save(out_path)
    print(f"✓ PPTX saved successfully at: {out_path}")

if __name__ == "__main__":
    create_deck()
